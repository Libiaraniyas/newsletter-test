#!/usr/bin/env python3
"""
generate_pptx.py — Benchmark Report PPTX Generator
Usage: python3 generate_pptx.py benchmark-Q1-2026-approved.json [output.pptx]
Requires: pip install python-pptx
"""

import json
import sys
import os
import textwrap
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Color palette (from template analysis) ─────────────────────
DARK_NAVY   = RGBColor(0x0E, 0x28, 0x41)
MID_BLUE    = RGBColor(0x15, 0x60, 0x82)
TEAL        = RGBColor(0x61, 0xC9, 0xA8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
NEAR_BLACK  = RGBColor(0x0A, 0x0A, 0x0A)
ORANGE      = RGBColor(0xE9, 0x71, 0x32)
LIGHT_TEAL  = RGBColor(0xD0, 0xF0, 0xE8)
DARK_BG     = RGBColor(0x12, 0x2B, 0x45)

# ── Slide dimensions — widescreen 13.33" × 7.5" ───────────────
SW = 12192000
SH = 6858000

COMPANIES = [
    "Unilever", "Mondelez", "PepsiCo", "Nestlé", "Danone",
    "Hershey", "Coca-Cola", "Kraft Heinz", "Keurig Dr Pepper",
    "General Mills", "Barry Callebaut", "Haier", "SharkNinja"
]
SOURCES = [
    "Earnings call transcripts",
    "Investor presentations",
    "10-Q / Annual Report filings",
    "Barclays analyst reports",
    "Jefferies analyst reports",
    "Goldman Sachs analyst reports",
]


# ── Shape helpers ───────────────────────────────────────────────

def set_bg(slide, color=DARK_NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill=None, line=False):
    shape = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if not line:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, size=20, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_para(tf, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    if text:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return p

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def truncate(text, max_chars=400):
    if not text:
        return ""
    text = str(text)
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rsplit(' ', 1)[0] + "…"


# ── Slide builders ──────────────────────────────────────────────

def slide_title(prs, quarter, year):
    """Slide 1 — Title"""
    slide = blank_slide(prs)
    set_bg(slide, DARK_NAVY)

    # Left teal accent stripe
    add_rect(slide, 0, 0, 380000, SH, fill=TEAL)

    # Main title
    txb = slide.shapes.add_textbox(Emu(700000), Emu(2000000), Emu(8500000), Emu(1500000))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Quarterly Benchmark Report"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Subtitle
    add_textbox(slide, 700000, 3600000, 8500000, 600000,
                f"{quarter} {year}  ·  Cross-Sector Macro & Strategic Analysis",
                size=26, color=TEAL)

    # Company count
    add_textbox(slide, 700000, 4300000, 8500000, 400000,
                "13 companies  ·  CPG & Appliances  ·  Real earnings call sources",
                size=18, color=LIGHT_GRAY)

    # Bottom teal line
    add_rect(slide, 700000, 5800000, 5000000, 60000, fill=TEAL)


def slide_scope(prs, quarter, year):
    """Slide 2 — Companies & Sources"""
    slide = blank_slide(prs)
    set_bg(slide, DARK_NAVY)

    # Title bar
    add_rect(slide, 0, 0, SW, 760000, fill=MID_BLUE)
    add_textbox(slide, 400000, 150000, 9000000, 500000,
                f"Scope of Analysis  —  {quarter} {year}",
                size=30, bold=True, color=WHITE)

    # Left block — Companies
    add_rect(slide, 300000, 900000, 5200000, 5700000, fill=DARK_BG)
    add_textbox(slide, 500000, 950000, 4800000, 500000,
                "Companies of Scope", size=22, bold=True, color=TEAL)

    txb = slide.shapes.add_textbox(Emu(500000), Emu(1550000), Emu(4800000), Emu(4800000))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for co in COMPANIES:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "  •  " + co
        run.font.size = Pt(20)
        run.font.color.rgb = WHITE

    # Right block — Sources
    add_rect(slide, 6300000, 900000, 5600000, 5700000, fill=DARK_BG)
    add_textbox(slide, 6500000, 950000, 5200000, 500000,
                "Sources", size=22, bold=True, color=TEAL)

    txb2 = slide.shapes.add_textbox(Emu(6500000), Emu(1550000), Emu(5200000), Emu(4800000))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for src in SOURCES:
        if first:
            p = tf2.paragraphs[0]
            first = False
        else:
            p = tf2.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "  •  " + src
        run.font.size = Pt(20)
        run.font.color.rgb = WHITE


def slide_section_divider(prs, label, subtitle=""):
    """Full-bleed dark section divider slide."""
    slide = blank_slide(prs)
    set_bg(slide, DARK_NAVY)

    # Teal left stripe
    add_rect(slide, 0, 0, 380000, SH, fill=TEAL)

    # Section label
    txb = slide.shapes.add_textbox(Emu(700000), Emu(2400000), Emu(10000000), Emu(1200000))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = WHITE

    if subtitle:
        add_textbox(slide, 700000, 3700000, 10000000, 500000,
                    subtitle, size=22, color=TEAL)

    add_rect(slide, 700000, 5800000, 4000000, 60000, fill=TEAL)


def slide_macro_overview(prs, themes):
    """Macro Themes overview — grid of theme cards (max 6)."""
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_GRAY)

    # Header
    add_rect(slide, 0, 0, SW, 760000, fill=DARK_NAVY)
    add_textbox(slide, 400000, 150000, 9000000, 500000,
                "Key Macro Environment Themes", size=30, bold=True, color=WHITE)

    themes = themes[:6]
    cols = 3
    rows = (len(themes) + cols - 1) // cols
    card_w = 3600000
    card_h = 2300000 if rows <= 2 else 1700000
    gap = 180000
    start_x = 300000
    start_y = 900000

    for i, t in enumerate(themes):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)

        # Card background
        add_rect(slide, x, y, card_w, card_h, fill=DARK_NAVY)
        # Teal header strip
        add_rect(slide, x, y, card_w, 350000, fill=TEAL)

        # Theme title (on teal strip)
        theme_text = truncate(t.get("theme", ""), 60)
        add_textbox(slide, x + 80000, y + 40000, card_w - 160000, 300000,
                    theme_text, size=15, bold=True, color=DARK_NAVY)

        # Mention count badge
        cnt = t.get("mention_count") or t.get("company_count") or 0
        add_textbox(slide, x + card_w - 480000, y + 380000, 400000, 300000,
                    str(cnt), size=28, bold=True, color=TEAL)

        # Summary text
        summary = truncate(t.get("summary", ""), 200)
        txb = slide.shapes.add_textbox(
            Emu(x + 80000), Emu(y + 380000), Emu(card_w - 180000), Emu(card_h - 430000))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = summary
        run.font.size = Pt(14)
        run.font.color.rgb = LIGHT_GRAY


def slide_macro_theme(prs, item):
    """One slide per macro executive_synthesis item — matches slides 6-8."""
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_GRAY)

    theme = item.get("theme", "")
    action = item.get("action", "")
    citation = item.get("citation", "")

    # Title (top-left, dark text)
    add_textbox(slide, 150000, 30000, 9000000, 700000,
                theme, size=32, bold=True, color=DARK_NAVY)

    # Teal rounded-rect bar — "What is happening"
    bar = add_rect(slide, 650000, 800000, 10785352, 620000, fill=TEAL)

    add_textbox(slide, 4500000, 840000, 4000000, 500000,
                "What is happening", size=22, bold=True, color=DARK_NAVY)

    # Body — action/summary
    body_text = action if action else citation
    txb = slide.shapes.add_textbox(
        Emu(900000), Emu(1530000), Emu(10400000), Emu(1400000))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = truncate(body_text, 350)
    run.font.size = Pt(20)
    run.font.color.rgb = NEAR_BLACK

    # Dark quote block (bottom)
    add_rect(slide, 650000, 3050000, 10785352, 3500000, fill=DARK_NAVY)

    # Citation inside quote block
    cite_text = truncate(citation, 500) if citation else ""
    txb2 = slide.shapes.add_textbox(
        Emu(900000), Emu(3200000), Emu(10200000), Emu(3200000))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = cite_text
    run2.font.size = Pt(18)
    run2.font.color.rgb = LIGHT_GRAY
    run2.font.italic = True


def slide_strategic_overview(prs, themes):
    """Strategic Themes overview — table layout."""
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_GRAY)

    # Header
    add_rect(slide, 0, 0, SW, 760000, fill=DARK_NAVY)
    add_textbox(slide, 400000, 150000, 9000000, 500000,
                "Key Strategic Themes", size=30, bold=True, color=WHITE)

    # Column headers
    col_x = [350000, 3200000, 4400000, 9200000]
    col_w = [2700000, 1000000, 4650000, 2800000]
    headers = ["Theme", "Companies", "Summary", "Tier"]
    row_y = 850000
    row_h = 380000

    add_rect(slide, 350000, row_y, SW - 700000, row_h, fill=MID_BLUE)
    for i, hdr in enumerate(headers):
        add_textbox(slide, col_x[i] + 60000, row_y + 60000, col_w[i], row_h - 80000,
                    hdr, size=16, bold=True, color=WHITE)

    row_y += row_h
    for idx, t in enumerate(themes[:12]):
        bg = DARK_BG if idx % 2 == 0 else DARK_NAVY
        add_rect(slide, 350000, row_y, SW - 700000, row_h, fill=bg)

        cnt = t.get("company_count") or t.get("mention_count") or 0
        tier = "Key Theme" if cnt >= 7 else "Supporting" if cnt >= 4 else "Outlier"
        tier_color = TEAL if cnt >= 7 else ORANGE if cnt >= 4 else LIGHT_GRAY

        cells = [
            (truncate(t.get("theme", ""), 55), WHITE, False),
            (str(cnt), TEAL, True),
            (truncate(t.get("summary", ""), 130), LIGHT_GRAY, False),
            (tier, tier_color, True),
        ]
        for i, (txt, clr, bld) in enumerate(cells):
            add_textbox(slide, col_x[i] + 60000, row_y + 50000, col_w[i] - 60000,
                        row_h - 60000, txt, size=13, bold=bld, color=clr)

        row_y += row_h
        if row_y + row_h > SH - 200000:
            break


def slide_deep_dive(prs, case):
    """One deep-dive slide — matches template slides 18-22."""
    slide = blank_slide(prs)
    set_bg(slide, LIGHT_GRAY)

    company      = case.get("company", "")
    init_type    = case.get("initiative_type", "")
    title        = case.get("title", "")
    announced    = truncate(case.get("what_announced", ""), 420)
    rationale    = truncate(case.get("strategic_rationale", ""), 420)
    financial    = truncate(case.get("financial_impact", ""), 400)
    exec_risks   = truncate(case.get("execution_risks", ""), 180)
    status       = case.get("execution_status", "announced")

    # Category title (top-left)
    add_textbox(slide, 150000, 0, 8000000, 650000,
                init_type, size=32, bold=True, color=DARK_NAVY)

    # Full-width teal topic bar
    add_rect(slide, 0, 700000, SW, 520000, fill=TEAL)
    add_textbox(slide, 100000, 730000, SW - 200000, 460000,
                f"  {company}  —  {title}", size=22, bold=True, color=DARK_NAVY)

    # Status chip (top-right on teal bar)
    status_colors = {
        "announced": ORANGE,
        "regulatory review": RGBColor(0x6B, 0x21, 0xA8),
        "in-progress": MID_BLUE,
        "completed": RGBColor(0x16, 0x6A, 0x24),
    }
    chip_color = status_colors.get(status, ORANGE)
    add_rect(slide, SW - 1900000, 740000, 1800000, 380000, fill=chip_color)
    add_textbox(slide, SW - 1850000, 760000, 1700000, 340000,
                status.upper(), size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Left section: upper block (What announced) ──────────────
    LBX, LBY, LBW, LBH = 600000, 1310000, 8130000, 2620000
    add_rect(slide, LBX, LBY, LBW, LBH, fill=DARK_NAVY)

    add_textbox(slide, LBX + 120000, LBY + 80000, 3500000, 400000,
                "What has been announced:", size=18, bold=True, color=TEAL)
    txb = slide.shapes.add_textbox(
        Emu(LBX + 120000), Emu(LBY + 520000), Emu(LBW - 200000), Emu(LBH - 600000))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = announced
    run.font.size = Pt(16)
    run.font.color.rgb = WHITE

    # ── Left section: lower block (Strategic Rationale) ─────────
    LBY2 = LBY + LBH + 110000
    LBH2 = SH - LBY2 - 100000
    add_rect(slide, LBX, LBY2, LBW, LBH2, fill=DARK_NAVY)

    add_textbox(slide, LBX + 120000, LBY2 + 80000, 3500000, 400000,
                "Strategic Rationale:", size=18, bold=True, color=TEAL)

    rat_text = rationale
    if exec_risks:
        rat_text += f"\n\nRisks: {exec_risks}"
    txb2 = slide.shapes.add_textbox(
        Emu(LBX + 120000), Emu(LBY2 + 520000), Emu(LBW - 200000), Emu(LBH2 - 600000))
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = rat_text
    run2.font.size = Pt(16)
    run2.font.color.rgb = WHITE

    # ── Right column: Financial Impact ───────────────────────────
    RCX = LBX + LBW + 110000
    RCW = SW - RCX - 100000
    add_rect(slide, RCX, LBY, RCW, SH - LBY - 100000, fill=MID_BLUE)

    add_textbox(slide, RCX + 100000, LBY + 80000, RCW - 150000, 400000,
                "Financial Impact", size=18, bold=True, color=WHITE)
    txb3 = slide.shapes.add_textbox(
        Emu(RCX + 100000), Emu(LBY + 540000), Emu(RCW - 150000),
        Emu(SH - LBY - 100000 - 600000))
    tf3 = txb3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = financial if financial else "Financial data not available from verified sources."
    run3.font.size = Pt(15)
    run3.font.color.rgb = WHITE


def slide_thank_you(prs):
    """Last slide — Thank You."""
    slide = blank_slide(prs)
    set_bg(slide, DARK_NAVY)

    add_rect(slide, 0, 0, 380000, SH, fill=TEAL)

    txb = slide.shapes.add_textbox(Emu(700000), Emu(2400000), Emu(10000000), Emu(1200000))
    tf = txb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Thank You"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = WHITE

    add_textbox(slide, 700000, 3700000, 8000000, 500000,
                "Sources verified via real earnings calls, filings, and analyst reports.",
                size=18, color=LIGHT_GRAY)

    add_rect(slide, 700000, 5800000, 4000000, 60000, fill=TEAL)


# ── Main ────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print("Usage: python3 generate_pptx.py benchmark-Q1-2026-approved.json [output.pptx]")
        sys.exit(1)

    json_path = sys.argv[1]
    if not os.path.exists(json_path):
        print(f"Error: file not found: {json_path}")
        sys.exit(1)

    with open(json_path, encoding='utf-8') as f:
        data = json.load(f)

    quarter = data.get("quarter", "Q?")
    year    = data.get("year", "????")
    macro   = data.get("macro", {})
    strat   = data.get("strategic", {})

    macro_synth  = macro.get("executive_synthesis", [])
    macro_themes = macro.get("theme_frequency_table", [])
    strat_themes = strat.get("theme_frequency_table", [])
    deep_dives   = strat.get("deep_dive_cases", [])

    out_path = sys.argv[2] if len(sys.argv) > 2 else \
        json_path.replace(".json", ".pptx").replace("-approved", "-report")

    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)

    # 1. Title
    slide_title(prs, quarter, year)

    # 2. Scope
    slide_scope(prs, quarter, year)

    # 3. Macro section divider
    slide_section_divider(prs,
        "Macro Environment Analysis",
        f"Consumer demand · Input costs · FX · Tariffs · Supply chain · Labor  —  {quarter} {year}")

    # 4. Macro themes overview (grid)
    if macro_themes:
        slide_macro_overview(prs, macro_themes)

    # 5. Individual macro theme slides (up to 5)
    for item in macro_synth[:5]:
        slide_macro_theme(prs, item)

    # 6. Strategic section divider
    slide_section_divider(prs,
        "Strategic Themes Analysis",
        f"M&A · Restructuring · Pricing pivots · Digital · Portfolio optimization  —  {quarter} {year}")

    # 7. Strategic themes overview (table)
    if strat_themes:
        slide_strategic_overview(prs, strat_themes)

    # 8. Deep Dives section divider
    if deep_dives:
        slide_section_divider(prs,
            "Deep Dives on Strategic Moves",
            f"{len(deep_dives)} company-level case studies")

    # 9. Deep dive slides
    for case in deep_dives:
        slide_deep_dive(prs, case)

    # 10. Thank You
    slide_thank_you(prs)

    prs.save(out_path)
    slide_count = len(prs.slides)
    print(f"Saved {slide_count} slides → {out_path}")


if __name__ == "__main__":
    main()
