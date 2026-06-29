import json
import base64
import io

def handler(event, context):
    cors = {
        "Access-Control-Allow-Origin": "https://news-digest-ag1.pages.dev",
        "Access-Control-Allow-Headers": "Content-Type",
        "Access-Control-Allow-Methods": "POST, OPTIONS",
    }

    if event.get("httpMethod") == "OPTIONS":
        return {"statusCode": 200, "headers": cors, "body": ""}
    if event.get("httpMethod") != "POST":
        return {"statusCode": 405, "headers": cors,
                "body": json.dumps({"error": "Method not allowed"})}

    try:
        data = json.loads(event.get("body") or "{}")
    except Exception:
        return {"statusCode": 400, "headers": cors,
                "body": json.dumps({"error": "Invalid JSON body"})}

    try:
        pptx_bytes = _generate(data)
    except Exception as e:
        return {"statusCode": 500, "headers": cors,
                "body": json.dumps({"error": str(e)})}

    quarter  = data.get("quarter", "Q?")
    year     = data.get("year", "????")
    filename = f"benchmark-{quarter}-{year}-report.pptx"
    encoded  = base64.b64encode(pptx_bytes).decode()

    return {
        "statusCode": 200,
        "headers": {**cors, "content-type": "application/json"},
        "body": json.dumps({"pptx_base64": encoded, "filename": filename}),
    }


# ── PPTX generation (inlined from generate_pptx.py) ─────────────

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_NAVY  = RGBColor(0x0E, 0x28, 0x41)
MID_BLUE   = RGBColor(0x15, 0x60, 0x82)
TEAL       = RGBColor(0x61, 0xC9, 0xA8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
NEAR_BLACK = RGBColor(0x0A, 0x0A, 0x0A)
ORANGE     = RGBColor(0xE9, 0x71, 0x32)
DARK_BG    = RGBColor(0x12, 0x2B, 0x45)

SW = 12192000
SH = 6858000

COMPANIES = [
    "Unilever", "Mondelez", "PepsiCo", "Nestlé", "Danone",
    "Hershey", "Coca-Cola", "Kraft Heinz", "Keurig Dr Pepper",
    "General Mills", "Barry Callebaut", "Haier", "SharkNinja",
]
SOURCES = [
    "Earnings call transcripts",
    "Investor presentations",
    "10-Q / Annual Report filings",
    "Barclays analyst reports",
    "Jefferies analyst reports",
    "Goldman Sachs analyst reports",
]


def _set_bg(slide, color=DARK_NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, fill=None):
    shape = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.line.fill.background()
    return shape


def _tb(slide, l, t, w, h, text, size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT):
    txb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text           = text
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.color.rgb = color
    return txb


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _trunc(text, n=400):
    if not text:
        return ""
    text = str(text)
    return text if len(text) <= n else text[:n].rsplit(" ", 1)[0] + "…"


def _slide_title(prs, quarter, year):
    s = _blank(prs)
    _set_bg(s, DARK_NAVY)
    _rect(s, 0, 0, 380000, SH, fill=TEAL)
    txb = s.shapes.add_textbox(Emu(700000), Emu(2000000), Emu(8500000), Emu(1500000))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    r   = p.add_run()
    r.text           = "Quarterly Benchmark Report"
    r.font.size      = Pt(48)
    r.font.bold      = True
    r.font.color.rgb = WHITE
    _tb(s, 700000, 3600000, 8500000, 600000,
        f"{quarter} {year}  ·  Cross-Sector Macro & Strategic Analysis",
        size=26, color=TEAL)
    _tb(s, 700000, 4300000, 8500000, 400000,
        "CPG & Appliances  ·  Real earnings call sources",
        size=18, color=LIGHT_GRAY)
    _rect(s, 700000, 5800000, 5000000, 60000, fill=TEAL)


def _slide_scope(prs, quarter, year):
    s = _blank(prs)
    _set_bg(s, DARK_NAVY)
    _rect(s, 0, 0, SW, 760000, fill=MID_BLUE)
    _tb(s, 400000, 150000, 9000000, 500000,
        f"Scope of Analysis  —  {quarter} {year}", size=30, bold=True)

    _rect(s, 300000, 900000, 5200000, 5700000, fill=DARK_BG)
    _tb(s, 500000, 950000, 4800000, 500000,
        "Companies of Scope", size=22, bold=True, color=TEAL)
    txb = s.shapes.add_textbox(Emu(500000), Emu(1550000), Emu(4800000), Emu(4800000))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for co in COMPANIES:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(6)
        r = p.add_run()
        r.text           = "  •  " + co
        r.font.size      = Pt(20)
        r.font.color.rgb = WHITE

    _rect(s, 6300000, 900000, 5600000, 5700000, fill=DARK_BG)
    _tb(s, 6500000, 950000, 5200000, 500000,
        "Sources", size=22, bold=True, color=TEAL)
    txb2 = s.shapes.add_textbox(Emu(6500000), Emu(1550000), Emu(5200000), Emu(4800000))
    tf2  = txb2.text_frame
    tf2.word_wrap = True
    first = True
    for src in SOURCES:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.space_after = Pt(8)
        r = p.add_run()
        r.text           = "  •  " + src
        r.font.size      = Pt(20)
        r.font.color.rgb = WHITE


def _slide_divider(prs, label, subtitle=""):
    s = _blank(prs)
    _set_bg(s, DARK_NAVY)
    _rect(s, 0, 0, 380000, SH, fill=TEAL)
    txb = s.shapes.add_textbox(Emu(700000), Emu(2400000), Emu(10000000), Emu(1200000))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    r   = p.add_run()
    r.text           = label
    r.font.size      = Pt(44)
    r.font.bold      = True
    r.font.color.rgb = WHITE
    if subtitle:
        _tb(s, 700000, 3700000, 10000000, 500000, subtitle, size=22, color=TEAL)
    _rect(s, 700000, 5800000, 4000000, 60000, fill=TEAL)


def _slide_macro_overview(prs, themes):
    s = _blank(prs)
    _set_bg(s, LIGHT_GRAY)
    _rect(s, 0, 0, SW, 760000, fill=DARK_NAVY)
    _tb(s, 400000, 150000, 9000000, 500000,
        "Key Macro Environment Themes", size=30, bold=True)

    themes = themes[:6]
    cols   = 3
    rows   = (len(themes) + cols - 1) // cols
    card_w = 3600000
    card_h = 2300000 if rows <= 2 else 1700000
    gap    = 180000
    sx, sy = 300000, 900000

    for i, t in enumerate(themes):
        col = i % cols
        row = i // cols
        x   = sx + col * (card_w + gap)
        y   = sy + row * (card_h + gap)
        _rect(s, x, y, card_w, card_h, fill=DARK_NAVY)
        _rect(s, x, y, card_w, 350000, fill=TEAL)
        _tb(s, x + 80000, y + 40000, card_w - 160000, 300000,
            _trunc(t.get("theme", ""), 60), size=15, bold=True, color=DARK_NAVY)
        cnt = t.get("mention_count") or t.get("company_count") or 0
        _tb(s, x + card_w - 480000, y + 380000, 400000, 300000,
            str(cnt), size=28, bold=True, color=TEAL)
        txb = s.shapes.add_textbox(Emu(x + 80000), Emu(y + 380000),
                                   Emu(card_w - 180000), Emu(card_h - 430000))
        tf  = txb.text_frame
        tf.word_wrap = True
        p   = tf.paragraphs[0]
        r   = p.add_run()
        r.text           = _trunc(t.get("summary", ""), 200)
        r.font.size      = Pt(14)
        r.font.color.rgb = LIGHT_GRAY


def _slide_macro_theme(prs, item):
    s = _blank(prs)
    _set_bg(s, LIGHT_GRAY)
    _tb(s, 150000, 30000, 9000000, 700000,
        item.get("theme", ""), size=32, bold=True, color=DARK_NAVY)
    _rect(s, 650000, 800000, 10785352, 620000, fill=TEAL)
    _tb(s, 4500000, 840000, 4000000, 500000,
        "What is happening", size=22, bold=True, color=DARK_NAVY)
    body = item.get("action") or item.get("citation", "")
    txb  = s.shapes.add_textbox(Emu(900000), Emu(1530000), Emu(10400000), Emu(1400000))
    tf   = txb.text_frame
    tf.word_wrap = True
    p    = tf.paragraphs[0]
    r    = p.add_run()
    r.text           = _trunc(body, 350)
    r.font.size      = Pt(20)
    r.font.color.rgb = NEAR_BLACK
    _rect(s, 650000, 3050000, 10785352, 3500000, fill=DARK_NAVY)
    cite = _trunc(item.get("citation", ""), 500)
    txb2 = s.shapes.add_textbox(Emu(900000), Emu(3200000), Emu(10200000), Emu(3200000))
    tf2  = txb2.text_frame
    tf2.word_wrap = True
    p2   = tf2.paragraphs[0]
    r2   = p2.add_run()
    r2.text           = cite
    r2.font.size      = Pt(18)
    r2.font.color.rgb = LIGHT_GRAY
    r2.font.italic    = True


def _slide_strategic_overview(prs, themes):
    s = _blank(prs)
    _set_bg(s, LIGHT_GRAY)
    _rect(s, 0, 0, SW, 760000, fill=DARK_NAVY)
    _tb(s, 400000, 150000, 9000000, 500000,
        "Key Strategic Themes", size=30, bold=True)

    col_x = [350000, 3200000, 4400000, 9200000]
    col_w = [2700000, 1000000, 4650000, 2800000]
    headers = ["Theme", "Companies", "Summary", "Tier"]
    ry, rh = 850000, 380000

    _rect(s, 350000, ry, SW - 700000, rh, fill=MID_BLUE)
    for i, hdr in enumerate(headers):
        _tb(s, col_x[i] + 60000, ry + 60000, col_w[i], rh - 80000,
            hdr, size=16, bold=True)
    ry += rh

    for idx, t in enumerate(themes[:12]):
        bg  = DARK_BG if idx % 2 == 0 else DARK_NAVY
        _rect(s, 350000, ry, SW - 700000, rh, fill=bg)
        cnt = t.get("company_count") or t.get("mention_count") or 0
        tier       = "Key Theme" if cnt >= 7 else "Supporting" if cnt >= 4 else "Outlier"
        tier_color = TEAL if cnt >= 7 else ORANGE if cnt >= 4 else LIGHT_GRAY
        cells = [
            (_trunc(t.get("theme", ""), 55), WHITE, False),
            (str(cnt),                        TEAL,  True),
            (_trunc(t.get("summary", ""), 130), LIGHT_GRAY, False),
            (tier,                             tier_color, True),
        ]
        for i, (txt, clr, bld) in enumerate(cells):
            _tb(s, col_x[i] + 60000, ry + 50000, col_w[i] - 60000, rh - 60000,
                txt, size=13, bold=bld, color=clr)
        ry += rh
        if ry + rh > SH - 200000:
            break


def _slide_deep_dive(prs, case):
    s = _blank(prs)
    _set_bg(s, LIGHT_GRAY)

    company   = case.get("company", "")
    init_type = case.get("initiative_type", "")
    title     = case.get("title", "")
    announced = _trunc(case.get("what_announced", ""), 420)
    rationale = _trunc(case.get("strategic_rationale", ""), 420)
    financial = _trunc(case.get("financial_impact", ""), 400)
    risks     = _trunc(case.get("execution_risks", ""), 180)
    status    = case.get("execution_status", "announced")

    _tb(s, 150000, 0, 8000000, 650000, init_type, size=32, bold=True, color=DARK_NAVY)
    _rect(s, 0, 700000, SW, 520000, fill=TEAL)
    _tb(s, 100000, 730000, SW - 200000, 460000,
        f"  {company}  —  {title}", size=22, bold=True, color=DARK_NAVY)

    status_colors = {
        "announced":        ORANGE,
        "regulatory review": RGBColor(0x6B, 0x21, 0xA8),
        "in-progress":      MID_BLUE,
        "completed":        RGBColor(0x16, 0x6A, 0x24),
    }
    chip = status_colors.get(status, ORANGE)
    _rect(s, SW - 1900000, 740000, 1800000, 380000, fill=chip)
    _tb(s, SW - 1850000, 760000, 1700000, 340000,
        status.upper(), size=14, bold=True, align=PP_ALIGN.CENTER)

    LBX, LBY, LBW, LBH = 600000, 1310000, 8130000, 2620000
    _rect(s, LBX, LBY, LBW, LBH, fill=DARK_NAVY)
    _tb(s, LBX + 120000, LBY + 80000, 3500000, 400000,
        "What has been announced:", size=18, bold=True, color=TEAL)
    txb = s.shapes.add_textbox(Emu(LBX + 120000), Emu(LBY + 520000),
                               Emu(LBW - 200000), Emu(LBH - 600000))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    r   = p.add_run()
    r.text           = announced
    r.font.size      = Pt(16)
    r.font.color.rgb = WHITE

    LBY2 = LBY + LBH + 110000
    LBH2 = SH - LBY2 - 100000
    _rect(s, LBX, LBY2, LBW, LBH2, fill=DARK_NAVY)
    _tb(s, LBX + 120000, LBY2 + 80000, 3500000, 400000,
        "Strategic Rationale:", size=18, bold=True, color=TEAL)
    rat = rationale + (f"\n\nRisks: {risks}" if risks else "")
    txb2 = s.shapes.add_textbox(Emu(LBX + 120000), Emu(LBY2 + 520000),
                                Emu(LBW - 200000), Emu(LBH2 - 600000))
    tf2  = txb2.text_frame
    tf2.word_wrap = True
    p2   = tf2.paragraphs[0]
    r2   = p2.add_run()
    r2.text           = rat
    r2.font.size      = Pt(16)
    r2.font.color.rgb = WHITE

    RCX = LBX + LBW + 110000
    RCW = SW - RCX - 100000
    _rect(s, RCX, LBY, RCW, SH - LBY - 100000, fill=MID_BLUE)
    _tb(s, RCX + 100000, LBY + 80000, RCW - 150000, 400000,
        "Financial Impact", size=18, bold=True)
    txb3 = s.shapes.add_textbox(Emu(RCX + 100000), Emu(LBY + 540000),
                                Emu(RCW - 150000), Emu(SH - LBY - 100000 - 600000))
    tf3  = txb3.text_frame
    tf3.word_wrap = True
    p3   = tf3.paragraphs[0]
    r3   = p3.add_run()
    r3.text           = financial or "Financial data not available from verified sources."
    r3.font.size      = Pt(15)
    r3.font.color.rgb = WHITE


def _slide_thank_you(prs):
    s = _blank(prs)
    _set_bg(s, DARK_NAVY)
    _rect(s, 0, 0, 380000, SH, fill=TEAL)
    txb = s.shapes.add_textbox(Emu(700000), Emu(2400000), Emu(10000000), Emu(1200000))
    tf  = txb.text_frame
    p   = tf.paragraphs[0]
    r   = p.add_run()
    r.text           = "Thank You"
    r.font.size      = Pt(54)
    r.font.bold      = True
    r.font.color.rgb = WHITE
    _tb(s, 700000, 3700000, 8000000, 500000,
        "Sources verified via real earnings calls, filings, and analyst reports.",
        size=18, color=LIGHT_GRAY)
    _rect(s, 700000, 5800000, 4000000, 60000, fill=TEAL)


def _generate(data):
    quarter      = data.get("quarter", "Q?")
    year         = data.get("year", "????")
    macro        = data.get("macro", {})
    strat        = data.get("strategic", {})
    macro_synth  = macro.get("executive_synthesis", [])
    macro_themes = macro.get("theme_frequency_table", [])
    strat_themes = strat.get("theme_frequency_table", [])
    deep_dives   = strat.get("deep_dive_cases", [])

    prs = Presentation()
    prs.slide_width  = Emu(SW)
    prs.slide_height = Emu(SH)

    _slide_title(prs, quarter, year)
    _slide_scope(prs, quarter, year)
    _slide_divider(prs, "Macro Environment Analysis",
                   f"Consumer demand · Input costs · FX · Tariffs · Supply chain · Labor  —  {quarter} {year}")
    if macro_themes:
        _slide_macro_overview(prs, macro_themes)
    for item in macro_synth[:5]:
        _slide_macro_theme(prs, item)
    _slide_divider(prs, "Strategic Themes Analysis",
                   f"M&A · Restructuring · Pricing pivots · Digital · Portfolio optimization  —  {quarter} {year}")
    if strat_themes:
        _slide_strategic_overview(prs, strat_themes)
    if deep_dives:
        _slide_divider(prs, "Deep Dives on Strategic Moves",
                       f"{len(deep_dives)} company-level case studies")
    for case in deep_dives:
        _slide_deep_dive(prs, case)
    _slide_thank_you(prs)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
